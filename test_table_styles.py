"""Generate a test PPTX with one slide per candidate table style GUID.

Each slide has a small sample table with the style applied so you can open
the file in PowerPoint and identify which GUID maps to which built-in style.

Usage:
    uv run python test_table_styles.py
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Candidate GUIDs — covers Dark Style 1 variants + common neighbours
STYLE_GUIDS = [
    ("{2D5ABB26-0587-4C30-8999-92F81FD0307C}", "Dark Style 1"),
    ("{3C2FFA5D-87B4-456A-9821-1D502468CF0F}", "Dark Style 1 - Accent 1"),
    ("{543D1CFD-3C2A-4F9E-B6E0-7FCA7B8E5C35}", "Candidate DS1-A1 alt"),
    ("{E929F9F4-4A8F-4326-A1B4-22849713DDAB}", "Candidate DS1 variant A"),
    ("{5202B0CA-FC54-4496-8BCA-5EF66A818D29}", "Dark Style 2"),
    ("{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}", "Medium Style 2 - Accent 1"),
    ("{68D230F3-CF80-4859-8CE7-A43EE81993B5}", "Candidate DS1 variant B"),
    ("{C4B1156A-380E-4F78-BDF5-A137D16BA258}", "Dark Style 1 - Accent 2"),
    ("{D7AC3CCA-C797-4891-BE02-D94E43425B78}", "Dark Style 1 - Accent 3"),
    ("{DF591B68-6B1E-4246-A7FC-B16DE0C081A6}", "Candidate DS1 variant C"),
    ("{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}", "Candidate DS1-A4 alt A"),
    ("{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}", "Dark Style 1 - Accent 5"),
    ("{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}", "Dark Style 1 - Accent 6"),
    ("{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}", "Dark Style 1 - Accent 4 alt"),
    ("{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}", "Candidate DS1 variant D"),
    ("{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}", "Dark Style 2 - Accent 1-2"),
    ("{B301B821-A1FF-4177-AEE7-76D212191A09}", "Medium Style 1 - Accent 4"),
]


def main() -> None:
    template_path = Path("pptx_template/template1.pptx")
    prs = Presentation(str(template_path))

    # Find the "Title Only" layout
    layout = None
    for sl in prs.slide_layouts:
        if sl.name == "Title Only":
            layout = sl
            break
    if layout is None:
        raise ValueError("Title Only layout not found")

    # Remove existing slides
    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst):
        rId = sld_id_lst[0].rId
        prs.part.drop_rel(rId)
        del sld_id_lst[0]

    for idx, (guid, label) in enumerate(STYLE_GUIDS, start=1):
        slide = prs.slides.add_slide(layout)

        # Set title
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text = f"Slide {idx}: {label}\n{guid}"
                break

        # Create sample table
        rows, cols = 5, 5
        shape = slide.shapes.add_table(
            rows, cols, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.5),
        )
        table = shape.table

        tbl_pr = table._tbl.tblPr
        tbl_pr.set("firstRow", "1")
        tbl_pr.set("lastRow", "1")
        tbl_pr.set("bandRow", "1")
        tbl_pr.set("bandCol", "0")
        tbl_pr.set("firstCol", "1")

        # Remove default tableStyleId and apply candidate GUID
        for child in list(tbl_pr):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("tableStyleId", "tblStyle"):
                tbl_pr.remove(child)
        style_el = etree.SubElement(tbl_pr, qn("a:tableStyleId"))
        style_el.text = guid

        # Populate sample data
        headers = ["App", "Forecast", "Spend", "Change $", "Change %"]
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
        for r in range(1, rows - 1):
            table.cell(r, 0).text = f"App {r}"
            table.cell(r, 1).text = f"${r * 100}K"
            table.cell(r, 2).text = f"${r * 110}K"
            table.cell(r, 3).text = f"${r * 10}K"
            table.cell(r, 4).text = f"+{r}.0%"
        # Total row
        table.cell(rows - 1, 0).text = "Total"
        table.cell(rows - 1, 1).text = "$600K"
        table.cell(rows - 1, 2).text = "$660K"
        table.cell(rows - 1, 3).text = "$60K"
        table.cell(rows - 1, 4).text = "+10.0%"

    out = Path("output/table_style_test.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    out.write_bytes(buf.read())
    print(f"Saved {len(STYLE_GUIDS)} style samples to {out}")
    for idx, (guid, label) in enumerate(STYLE_GUIDS, start=1):
        print(f"  Slide {idx}: {label} — {guid}")


if __name__ == "__main__":
    main()
