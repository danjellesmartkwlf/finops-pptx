"""Config-driven PowerPoint slide builder for the FinOps Report Generator.

Reads slide definitions from ``slides_config.yaml`` and dispatches each
entry to a typed handler.  All formatting and layout helpers live in
``pptx_utils.py``.
"""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.pptx_utils import (
    add_paginated_table_slides,
    apply_table_style,
    build_table_on_slide,
    compute_mom_headers,
    fmt_abbreviated,
    fmt_abbreviated_signed,
    fmt_abbreviated_signed_whole,
    fmt_abbreviated_whole,
    fmt_daily_rate,
    fmt_daily_rate_signed,
    fmt_obs_abbreviated,
    fmt_pct,
    fmt_unit_cost,
    get_layout_by_name,
    remove_all_slides,
    resolve_template_path,
    set_cell,
    set_text_by_idx,
)

# ---------------------------------------------------------------------------
# KPI grid constants
# ---------------------------------------------------------------------------

# Background color palette: [row][col] = RGBColor
# Row 0 = AWN (blue), Row 1 = Cylance (teal), Row 2 = Totals (charcoal)
# Col 2 (Total column) uses a slightly darker shade of its row color.
_KPI_COLORS: list[list[RGBColor]] = [
    [RGBColor(0x13, 0x55, 0xAA), RGBColor(0x17, 0x65, 0xBB), RGBColor(0x0C, 0x3F, 0x88)],
    [RGBColor(0x00, 0x72, 0x64), RGBColor(0x00, 0x83, 0x74), RGBColor(0x00, 0x55, 0x4A)],
    [RGBColor(0x34, 0x48, 0x54), RGBColor(0x42, 0x59, 0x66), RGBColor(0x1E, 0x2B, 0x33)],
]

_KPI_ROW_LABELS = ["AWN", "Cylance", "Total"]
_KPI_COL_LABELS = ["AWS", "Databricks", "Total"]
_KPI_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_KPI_GREEN = RGBColor(0x00, 0xE6, 0x76)   # favorable (cost down)
_KPI_ORANGE = RGBColor(0xFF, 0x9A, 0x3C)  # unfavorable (cost up)

# Colors for the executive summary table
_COLOR_YELLOW = RGBColor(0xFF, 0xFF, 0x00)   # previous month values
_COLOR_GREEN = RGBColor(0x00, 0xB0, 0x50)    # favorable change (cost down)
_COLOR_ORANGE = RGBColor(0xFF, 0x73, 0x0E)   # unfavorable change (cost up)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide config loading
# ---------------------------------------------------------------------------

_SLIDES_CONFIG: dict[str, Any] | None = None


def _load_slides_config(project_root: Path | None = None) -> dict[str, Any]:
    """Load and cache slides_config.yaml."""
    global _SLIDES_CONFIG
    if _SLIDES_CONFIG is not None:
        return _SLIDES_CONFIG

    if project_root is None:
        project_root = Path(__file__).resolve().parent.parent
    config_path = project_root / "slides_config.yaml"
    with open(config_path, "r") as f:
        _SLIDES_CONFIG = yaml.safe_load(f)
    return _SLIDES_CONFIG


# ---------------------------------------------------------------------------
# Formatter registry
# ---------------------------------------------------------------------------

_FORMATTER_MAP: dict[str, Any] = {
    "abbreviated_currency": fmt_abbreviated,
    "abbreviated_currency_whole": fmt_abbreviated_whole,
    "signed_currency_whole": fmt_abbreviated_signed_whole,
    "signed_percentage": fmt_pct,
    "unit_cost": fmt_unit_cost,
    "obs_abbreviated": fmt_obs_abbreviated,
    "raw": None,
}


def _format_value(value: Any, format_name: str | None) -> str:
    """Apply a named formatter to a value."""
    if format_name is None or format_name == "raw":
        return str(value) if value is not None else ""
    func = _FORMATTER_MAP.get(format_name)
    if func is None:
        return str(value) if value is not None else ""
    return func(value)


# ---------------------------------------------------------------------------
# Data resolution helpers
# ---------------------------------------------------------------------------

def _resolve_data_path(data_source: str, context: dict[str, Any]) -> Any:
    """Resolve a dotted path like 'app_data.category_rollup' from context."""
    parts = data_source.split(".")
    obj = context
    for part in parts:
        if isinstance(obj, dict):
            obj = obj.get(part)
        else:
            return None
        if obj is None:
            return None
    return obj


def _check_requires(requires: str | None, context: dict[str, Any]) -> bool:
    """Return True if the ``requires`` dotted path is truthy (or absent)."""
    if requires is None:
        return True
    return bool(_resolve_data_path(requires, context))


def _resolve_headers(
    headers: list[str],
    prev_month: str,
    curr_month: str,
) -> list[str]:
    """Replace ``{prev_month}`` and ``{curr_month}`` tokens in headers."""
    return [
        h.replace("{prev_month}", prev_month).replace("{curr_month}", curr_month)
        for h in headers
    ]


# ---------------------------------------------------------------------------
# Row builders
# ---------------------------------------------------------------------------

def _build_data_rows(
    data: list[dict],
    columns: list[dict],
) -> list[list[str]]:
    """Build formatted table rows from data dicts and column config."""
    rows = []
    for record in data:
        row = []
        for col in columns:
            field = col["field"]
            fmt_name = col.get("format")
            fallback_field = col.get("fallback_field")

            value = record.get(field)

            # fallback_field convention: if the record has a non-None
            # fallback value, use that raw string instead of formatting.
            if fallback_field:
                fallback_val = record.get(fallback_field)
                if fallback_val:
                    row.append(fallback_val)
                    continue

            row.append(_format_value(value, fmt_name))
        rows.append(row)
    return rows


def _build_total_row(
    totals: dict,
    totals_columns: list[dict],
) -> list[str]:
    """Build a formatted total row from totals dict and column config."""
    row = []
    for col in totals_columns:
        if col.get("field") is None:
            row.append(col.get("text", "Total"))
        else:
            row.append(_format_value(totals.get(col["field"]), col.get("format")))
    return row


# ---------------------------------------------------------------------------
# Slide handlers
# ---------------------------------------------------------------------------

def _handle_title(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render the title slide."""
    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return
    slide = prs.slides.add_slide(layout)
    fmt = context["fmt"]
    set_text_by_idx(slide, 0, slide_cfg.get("heading", "").format(**fmt))
    set_text_by_idx(slide, 10, slide_cfg.get("subtitle", "").format(**fmt))


def _handle_transition(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a transition (section divider) slide."""
    if not _check_requires(slide_cfg.get("requires"), context):
        return
    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return
    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))


def _handle_content(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a content / narrative slide."""
    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return
    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))

    bucket = slide_cfg.get("bucket", "")
    narratives = context.get("narratives", {})
    set_text_by_idx(slide, 11, narratives.get(bucket, ""))


def _handle_chart(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a chart image slide.  Skipped when no chart for the bucket."""
    chart_images = context.get("chart_images")
    bucket = slide_cfg.get("bucket", "")
    if not chart_images or bucket not in chart_images:
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return
    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))

    png_stream = BytesIO(chart_images[bucket])
    slide.shapes.add_picture(
        png_stream,
        left=Inches(0.53),
        top=Inches(1.40),
        width=Inches(12.27),
        height=Inches(5.80),
    )


def _handle_table(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a paginated table slide from YAML config."""
    if not _check_requires(slide_cfg.get("requires"), context):
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return

    data = _resolve_data_path(slide_cfg["data_source"], context)
    if data is None:
        return

    totals_source = slide_cfg.get("totals_source")
    totals = _resolve_data_path(totals_source, context) if totals_source else None

    headers = _resolve_headers(
        slide_cfg["headers"],
        context.get("prev_month", ""),
        context.get("curr_month", ""),
    )

    data_rows = _build_data_rows(data, slide_cfg["columns"])

    total_row = None
    if totals and "totals_columns" in slide_cfg:
        total_row = _build_total_row(totals, slide_cfg["totals_columns"])

    add_paginated_table_slides(
        prs,
        layout,
        slide_cfg["title"],
        headers,
        data_rows,
        total_row=total_row,
        col_widths=slide_cfg.get("col_widths"),
        font_size=slide_cfg.get("font_size"),  # None → auto-compute
    )


def _handle_split_table(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render two tables stacked vertically on one slide.

    The lower table's vertical position is computed dynamically from the
    upper table's rendered height plus ``gap_below``.
    """
    if not _check_requires(slide_cfg.get("requires"), context):
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return
    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))

    tables_cfg = slide_cfg["tables"]
    upper_bottom = 1.4  # tracks bottom edge (inches) for lower table placement

    for tbl_cfg in tables_cfg:
        data = _resolve_data_path(tbl_cfg["data_source"], context)
        if data is None:
            data = []

        data_rows = _build_data_rows(data, tbl_cfg["columns"])
        if not data_rows:
            continue

        headers = _resolve_headers(
            tbl_cfg["headers"],
            context.get("prev_month", ""),
            context.get("curr_month", ""),
        )

        if tbl_cfg["position"] == "upper":
            tbl_top = tbl_cfg.get("top", 1.4)
        else:
            gap = tables_cfg[0].get("gap_below", 0.8)
            tbl_top = upper_bottom + gap

        _, bottom = build_table_on_slide(
            slide,
            headers,
            data_rows,
            col_widths=tbl_cfg.get("col_widths", [3.0, 2.3, 2.3, 2.3, 2.4]),
            font_size=tbl_cfg.get("font_size", 18),
            top_inches=tbl_top,
            row_height_factor=tbl_cfg.get("row_height_factor", 0.40),
        )

        if tbl_cfg["position"] == "upper":
            upper_bottom = bottom


def _handle_history_table(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a multi-month history table with colored MOM $ and MOM % columns.

    Expects rows in ``data_source`` with keys ``app``, ``monthly_values``
    (list of floats), ``mom_change``, and ``mom_pct``.  Month column headers
    come from ``month_labels_source``.
    """
    if not _check_requires(slide_cfg.get("requires"), context):
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return

    rows_data = _resolve_data_path(slide_cfg["data_source"], context)
    if not rows_data:
        return

    month_labels = _resolve_data_path(slide_cfg["month_labels_source"], context) or []
    totals = (
        _resolve_data_path(slide_cfg["totals_source"], context)
        if slide_cfg.get("totals_source")
        else None
    )

    col_widths: list[float] = slide_cfg.get(
        "col_widths", [2.0, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.55, 1.55]
    )
    font_size: int = slide_cfg.get("font_size", 13)

    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))

    headers = ["APPLICATION"] + list(month_labels) + ["MOM $", "MOM %"]
    n_cols = len(headers)
    n_data_rows = len(rows_data)
    n_rows = 1 + n_data_rows + (1 if totals else 0)
    col_widths = col_widths[:n_cols]

    shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(min(0.40 * n_rows, 5.8)),
    )
    table = shape.table

    tbl_pr = table._tbl.tblPr
    tbl_pr.set("firstRow", "1")
    tbl_pr.set("lastRow", "1" if totals else "0")
    tbl_pr.set("bandRow", "1")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstCol", "1")
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for col_idx, h in enumerate(headers):
        align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        set_cell(table.cell(0, col_idx), h, font_size=font_size, bold=True, alignment=align)

    # Data rows
    for row_idx, row in enumerate(rows_data, start=1):
        monthly_values: list[float] = row.get("monthly_values", [])
        mom_change: float = row.get("mom_change", 0.0)
        mom_pct: float = row.get("mom_pct", 0.0)
        mom_color = _COLOR_ORANGE if mom_change > 0 else _COLOR_GREEN

        set_cell(
            table.cell(row_idx, 0),
            row.get("app", ""),
            font_size=font_size,
            alignment=PP_ALIGN.LEFT,
        )
        for mi, val in enumerate(monthly_values):
            set_cell(
                table.cell(row_idx, 1 + mi),
                fmt_abbreviated(val),
                font_size=font_size,
                alignment=PP_ALIGN.RIGHT,
            )
        mom_col = 1 + len(monthly_values)
        set_cell(
            table.cell(row_idx, mom_col),
            fmt_abbreviated_signed(mom_change),
            font_size=font_size,
            alignment=PP_ALIGN.RIGHT,
            color=mom_color,
        )
        set_cell(
            table.cell(row_idx, mom_col + 1),
            fmt_pct(mom_pct),
            font_size=font_size,
            alignment=PP_ALIGN.RIGHT,
            color=mom_color,
        )

    # Totals row
    if totals:
        t_idx = 1 + n_data_rows
        monthly_values = totals.get("monthly_values", [])
        mom_change = totals.get("mom_change", 0.0)
        mom_pct = totals.get("mom_pct", 0.0)
        mom_color = _COLOR_ORANGE if mom_change > 0 else _COLOR_GREEN

        set_cell(
            table.cell(t_idx, 0), "Total", font_size=font_size, bold=True, alignment=PP_ALIGN.LEFT
        )
        for mi, val in enumerate(monthly_values):
            set_cell(
                table.cell(t_idx, 1 + mi),
                fmt_abbreviated(val),
                font_size=font_size,
                bold=True,
                alignment=PP_ALIGN.RIGHT,
            )
        mom_col = 1 + len(monthly_values)
        set_cell(
            table.cell(t_idx, mom_col),
            fmt_abbreviated_signed(mom_change),
            font_size=font_size,
            bold=True,
            alignment=PP_ALIGN.RIGHT,
            color=mom_color,
        )
        set_cell(
            table.cell(t_idx, mom_col + 1),
            fmt_pct(mom_pct),
            font_size=font_size,
            bold=True,
            alignment=PP_ALIGN.RIGHT,
            color=mom_color,
        )

    apply_table_style(table)


def _handle_summary_table(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render the executive summary MoM table (COGS/OPEX/Grand Total + daily rates).

    Expects ``context["summary_data"]`` to be populated by generate_report.py.
    Columns: Metric | <two months ago> | <previous month (yellow)> | Change ($) | Change (%)
    """
    data_key = slide_cfg.get("data_key", "summary_data")
    summary_data = context.get(data_key)
    if not summary_data:
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return

    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, slide_cfg.get("title", ""))

    col_two = summary_data["col_older"]
    col_prev = summary_data["col_newer"]
    headers = ["Metric", col_two, col_prev, "Change ($)", "Change (%)"]
    col_widths = slide_cfg.get("col_widths", [3.0, 2.3, 2.3, 2.3, 2.4])
    font_size = slide_cfg.get("font_size", 20)
    whole_numbers = slide_cfg.get("whole_numbers", False)
    _fmt_currency = fmt_abbreviated_whole if whole_numbers else fmt_abbreviated
    _fmt_currency_signed = fmt_abbreviated_signed_whole if whole_numbers else fmt_abbreviated_signed

    rows_data = summary_data["rows"]
    n_rows = 1 + len(rows_data)

    shape = slide.shapes.add_table(
        n_rows, 5,
        Inches(0.5), Inches(1.4), Inches(12.3),
        Inches(min(0.65 * n_rows, 5.8)),
    )
    table = shape.table

    tbl_pr = table._tbl.tblPr
    tbl_pr.set("firstRow", "1")
    tbl_pr.set("lastRow", "0")
    tbl_pr.set("bandRow", "1")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstCol", "1")
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for col_idx, h in enumerate(headers):
        align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        set_cell(table.cell(0, col_idx), h, font_size=font_size, bold=True, alignment=align)

    # Data rows
    for row_idx, row in enumerate(rows_data, start=1):
        is_bold = row.get("is_bold", False)
        row_type = row.get("row_type", "currency")
        change_val = row["change_dollar"]

        if row_type == "daily_rate":
            two_fmt = fmt_daily_rate(row["older"])
            prev_fmt = fmt_daily_rate(row["newer"])
            change_dollar_fmt = fmt_daily_rate_signed(change_val)
        else:
            two_fmt = _fmt_currency(row["older"])
            prev_fmt = _fmt_currency(row["newer"])
            change_dollar_fmt = _fmt_currency_signed(change_val)

        change_pct_fmt = fmt_pct(row["change_pct"])
        change_color = _COLOR_ORANGE if change_val > 0 else _COLOR_GREEN

        set_cell(table.cell(row_idx, 0), row["metric"], font_size=font_size, bold=is_bold, alignment=PP_ALIGN.LEFT)
        set_cell(table.cell(row_idx, 1), two_fmt, font_size=font_size, bold=is_bold, alignment=PP_ALIGN.RIGHT)
        set_cell(table.cell(row_idx, 2), prev_fmt, font_size=font_size, bold=is_bold, alignment=PP_ALIGN.RIGHT, color=_COLOR_YELLOW)
        set_cell(table.cell(row_idx, 3), change_dollar_fmt, font_size=font_size, bold=is_bold, alignment=PP_ALIGN.RIGHT, color=change_color)
        set_cell(table.cell(row_idx, 4), change_pct_fmt, font_size=font_size, bold=is_bold, alignment=PP_ALIGN.RIGHT, color=change_color)

    apply_table_style(table)


# ---------------------------------------------------------------------------
# KPI grid slide renderer
# ---------------------------------------------------------------------------

def _render_kpi_grid(slide: Any, grid: list[list[dict]]) -> None:
    """Draw a 3×3 KPI card grid onto *slide*.

    *grid* is a 3-row × 3-col list of cell dicts with keys:
    ``value``, ``prev``, ``change``, ``pct``.
    """
    left_margin = 0.4
    top_margin = 1.3
    total_w = 12.53
    total_h = 5.9
    gap = 0.12
    card_w = (total_w - 2 * gap) / 3
    card_h = (total_h - 2 * gap) / 3

    for r in range(3):
        for c in range(3):
            cell = grid[r][c]
            left = left_margin + c * (card_w + gap)
            top = top_margin + r * (card_h + gap)
            color = _KPI_COLORS[r][c]

            value = cell["value"]
            change = cell["change"]
            pct = cell["pct"]

            # Background card (filled rectangle, no visible border)
            shape = slide.shapes.add_shape(
                1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(card_w), Inches(card_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.color.rgb = color  # border matches fill → invisible

            tf = shape.text_frame
            tf.word_wrap = False
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.05)
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.08)

            # Line 1: "AWN  ·  AWS" label (small, top-left)
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.LEFT
            run0 = p0.add_run()
            run0.text = f"{_KPI_ROW_LABELS[r]}  ·  {_KPI_COL_LABELS[c]}"
            run0.font.size = Pt(11)
            run0.font.bold = False
            run0.font.color.rgb = _KPI_WHITE

            # Line 2: main dollar value (large, centered)
            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(4)
            run1 = p1.add_run()
            run1.text = fmt_abbreviated(value)
            run1.font.size = Pt(28)
            run1.font.bold = True
            run1.font.color.rgb = _KPI_WHITE

            # Line 3: MoM delta (small, centered, color-coded)
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(3)
            run2 = p2.add_run()
            sign_char = "▲" if change > 0 else ("▼" if change < 0 else "–")
            run2.text = f"{sign_char} {fmt_abbreviated_signed(change)}  {fmt_pct(pct)}"
            run2.font.size = Pt(12)
            run2.font.bold = False
            run2.font.color.rgb = _KPI_ORANGE if change > 0 else _KPI_GREEN


def _handle_kpi_grid(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render a 3×3 KPI cost overview grid slide."""
    if not _check_requires(slide_cfg.get("requires"), context):
        return

    kpi_grid_data = context.get("kpi_grid_data")
    if not kpi_grid_data:
        return

    bucket = slide_cfg.get("bucket", "Total")
    grid = kpi_grid_data.get(bucket)
    if not grid:
        return

    layout = layouts.get(slide_cfg["layout"])
    if layout is None:
        return

    slide = prs.slides.add_slide(layout)
    fmt = context.get("fmt", {})
    title = slide_cfg.get("title", "").format(**fmt)
    set_text_by_idx(slide, 0, title)
    _render_kpi_grid(slide, grid)


# ---------------------------------------------------------------------------
# COGS Forecast Variance slide
# ---------------------------------------------------------------------------

def _fmt_variance_pct(var_change: float | None, forecast: float | None) -> str:
    """Return variance % string, using 'n/m' when forecast is near zero."""
    if var_change is None or forecast is None:
        return "N/A"
    if abs(forecast) < 500:
        return "n/m"
    pct = (var_change / forecast) * 100.0
    sign = "+" if pct > 0 else ""
    return f"{sign}{pct:.1f}%"


def _handle_cogs_forecast_variance(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Render COGS Forecast Variance slide grouped by category with subtotals.

    Columns: Application/Category | Actual | Forecast | Variance ($) | Variance (%)
    Orange = over forecast (unfavorable), green = under forecast (favorable).
    Category subtotal rows are inserted after each category's apps.
    """
    if not _check_requires(slide_cfg.get("requires"), context):
        return

    app_metrics: list[dict] | None = _resolve_data_path("app_data.app_metrics", context)
    app_totals: dict | None = _resolve_data_path("app_data.app_totals", context)
    if not app_metrics:
        return

    fmt = context.get("fmt", {})
    month = fmt.get("month", "")
    year = fmt.get("year", "")
    title = f"COGS Forecast Variance \u2014 {month} {year}"

    layout = layouts.get(slide_cfg.get("layout", "title_only"))
    if layout is None:
        return

    # Group apps by category (preserve descending-spend order within category)
    from collections import defaultdict
    category_apps: dict[str, list[dict]] = defaultdict(list)
    for m in app_metrics:
        category_apps[m["category"]].append(m)

    # Sort categories by total current_month spend descending
    category_order = sorted(
        category_apps.keys(),
        key=lambda c: sum(m["current_month"] for m in category_apps[c]),
        reverse=True,
    )

    # Build flat row list: app rows + category subtotal row per category
    TableRow = dict  # type alias for clarity
    flat_rows: list[TableRow] = []
    for cat in category_order:
        apps = category_apps[cat]
        for m in apps:
            flat_rows.append({"_type": "app", **m})
        # Category subtotal
        cat_actual = sum(m["current_month"] for m in apps)
        cat_forecast = sum(m["forecast"] for m in apps if m["forecast"] is not None)
        cat_var = cat_actual - cat_forecast if any(m["forecast"] is not None for m in apps) else None
        flat_rows.append({
            "_type": "subtotal",
            "app": f"{cat} Category Total",
            "current_month": cat_actual,
            "forecast": cat_forecast if any(m["forecast"] is not None for m in apps) else None,
            "var_change": cat_var,
        })

    # Grand total row
    grand_actual = app_totals["current_month"] if app_totals else sum(m["current_month"] for m in app_metrics)
    grand_forecast = app_totals["forecast"] if app_totals else None
    grand_var = app_totals["var_change"] if app_totals else None

    font_size = slide_cfg.get("font_size", 16)
    col_widths = slide_cfg.get("col_widths", [3.0, 2.3, 2.3, 2.3, 2.4])
    headers = ["Application /\nCategory", f"{month} Actual", f"{month} Forecast", "Variance ($)", "Variance (%)"]

    n_data_rows = len(flat_rows)
    n_rows = 1 + n_data_rows + 1  # header + data + grand total

    slide = prs.slides.add_slide(layout)
    set_text_by_idx(slide, 0, title)

    shape = slide.shapes.add_table(
        n_rows, 5,
        Inches(0.5), Inches(1.4), Inches(12.3),
        Inches(min(0.32 * n_rows, 5.8)),
    )
    table = shape.table

    tbl_pr = table._tbl.tblPr
    tbl_pr.set("firstRow", "1")
    tbl_pr.set("lastRow", "1")
    tbl_pr.set("bandRow", "1")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstCol", "1")

    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for col_idx, h in enumerate(headers):
        align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        set_cell(table.cell(0, col_idx), h, font_size=font_size, bold=True, alignment=align)

    # Data rows
    for row_idx, row in enumerate(flat_rows, start=1):
        is_subtotal = row["_type"] == "subtotal"
        var_change = row.get("var_change")
        forecast = row.get("forecast")
        actual = row.get("current_month", 0.0)

        var_color = _COLOR_ORANGE if (var_change or 0) > 0 else _COLOR_GREEN
        var_dollar_str = fmt_abbreviated_signed(var_change) if var_change is not None else "N/A"
        var_pct_str = _fmt_variance_pct(var_change, forecast)

        set_cell(
            table.cell(row_idx, 0), row.get("app", ""),
            font_size=font_size, bold=is_subtotal, alignment=PP_ALIGN.LEFT,
        )
        set_cell(
            table.cell(row_idx, 1), fmt_abbreviated(actual),
            font_size=font_size, bold=is_subtotal, alignment=PP_ALIGN.RIGHT,
        )
        set_cell(
            table.cell(row_idx, 2), fmt_abbreviated(forecast) if forecast else "N/A",
            font_size=font_size, bold=is_subtotal, alignment=PP_ALIGN.RIGHT,
        )
        set_cell(
            table.cell(row_idx, 3), var_dollar_str,
            font_size=font_size, bold=is_subtotal, alignment=PP_ALIGN.RIGHT,
            color=var_color if var_change is not None else None,
        )
        set_cell(
            table.cell(row_idx, 4), var_pct_str,
            font_size=font_size, bold=is_subtotal, alignment=PP_ALIGN.RIGHT,
            color=var_color if var_change is not None else None,
        )

    # Grand total row
    grand_row_idx = 1 + n_data_rows
    grand_var_color = _COLOR_ORANGE if (grand_var or 0) > 0 else _COLOR_GREEN
    set_cell(table.cell(grand_row_idx, 0), "COGS Grand Total", font_size=font_size, bold=True, alignment=PP_ALIGN.LEFT)
    set_cell(table.cell(grand_row_idx, 1), fmt_abbreviated(grand_actual), font_size=font_size, bold=True, alignment=PP_ALIGN.RIGHT)
    set_cell(table.cell(grand_row_idx, 2), fmt_abbreviated(grand_forecast) if grand_forecast else "N/A", font_size=font_size, bold=True, alignment=PP_ALIGN.RIGHT)
    set_cell(table.cell(grand_row_idx, 3), fmt_abbreviated_signed(grand_var) if grand_var is not None else "N/A", font_size=font_size, bold=True, alignment=PP_ALIGN.RIGHT, color=grand_var_color if grand_var is not None else None)
    set_cell(table.cell(grand_row_idx, 4), _fmt_variance_pct(grand_var, grand_forecast), font_size=font_size, bold=True, alignment=PP_ALIGN.RIGHT, color=grand_var_color if grand_var is not None else None)

    apply_table_style(table)


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_HANDLER_MAP: dict[str, Any] = {
    "title": _handle_title,
    "transition": _handle_transition,
    "content": _handle_content,
    "chart": _handle_chart,
    "table": _handle_table,
    "split_table": _handle_split_table,
    "summary_table": _handle_summary_table,
    "history_table": _handle_history_table,
    "kpi_grid": _handle_kpi_grid,
    "cogs_forecast_variance": _handle_cogs_forecast_variance,
}


def _dispatch_slide(
    prs: Any,
    slide_cfg: dict[str, Any],
    layouts: dict[str, Any],
    context: dict[str, Any],
) -> None:
    """Dispatch a single slide config entry to its handler."""
    slide_type = slide_cfg.get("type", "")
    handler = _HANDLER_MAP.get(slide_type)
    if handler is None:
        logger.warning("Unknown slide type '%s' -- skipping.", slide_type)
        return
    handler(prs, slide_cfg, layouts, context)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_pptx(
    narratives: dict[str, str],
    all_metrics: list[dict],
    config: dict[str, Any],
    month_label: str = "",
    project_root: Path | None = None,
    chart_images: dict[str, bytes] | None = None,
    app_data: dict[str, Any] | None = None,
    summary_data: dict[str, Any] | None = None,
    cylance_summary_data: dict[str, Any] | None = None,
    cylance_dbx_summary_data: dict[str, Any] | None = None,
    unit_cost_data: dict[str, Any] | None = None,
    data_platform_data: dict[str, Any] | None = None,
    dbx_data: dict[str, Any] | None = None,
    dbx_summary_data: dict[str, Any] | None = None,
    data_lake_data: dict[str, Any] | None = None,
    kpi_grid_data: dict[str, Any] | None = None,
    other_app_data: dict[str, Any] | None = None,
) -> bytes:
    """Build a PowerPoint deck from ``slides_config.yaml`` definitions.

    Args:
        narratives: Bucket name -> narrative text.
        all_metrics: Metric dicts (reserved for future use).
        config: Parsed config.yaml (used for template resolution).
        month_label: e.g. ``"January 2026"`` for title formatting.
        project_root: Optional project root for template resolution.
        chart_images: Optional dict mapping bucket name to PNG bytes.
        app_data: Optional Phase 2/3 data dict.

    Returns:
        The ``.pptx`` file as raw bytes.
    """
    slides_config = _load_slides_config(project_root=project_root)
    template_path = resolve_template_path(config, project_root=project_root)
    prs = Presentation(str(template_path))

    # Resolve layout names → layout objects
    layout_map: dict[str, Any] = {}
    for key, name in slides_config.get("layouts", {}).items():
        try:
            layout_map[key] = get_layout_by_name(prs, name)
        except ValueError:
            logger.warning("Layout '%s' (%s) not found.", name, key)
            layout_map[key] = None

    remove_all_slides(prs)

    # Build format dict from month_label (e.g. "January 2026")
    parts = month_label.split() if month_label else []
    fmt = {
        "month": parts[0] if parts else "",
        "year": parts[1] if len(parts) > 1 else "",
    }

    # Compute MoM header labels for template token resolution
    prev_month = ""
    curr_month = ""
    if month_label:
        prev_month, curr_month = compute_mom_headers(month_label)

    # Rendering context available to all handlers
    context: dict[str, Any] = {
        "config": config,
        "narratives": narratives,
        "all_metrics": all_metrics,
        "month_label": month_label,
        "chart_images": chart_images,
        "app_data": app_data,
        "summary_data": summary_data,
        "cylance_summary_data": cylance_summary_data,
        "cylance_dbx_summary_data": cylance_dbx_summary_data,
        "unit_cost_data": unit_cost_data,
        "data_platform_data": data_platform_data,
        "dbx_data": dbx_data,
        "dbx_summary_data": dbx_summary_data,
        "data_lake_data": data_lake_data,
        "kpi_grid_data": kpi_grid_data,
        "other_app_data": other_app_data,
        "fmt": fmt,
        "prev_month": prev_month,
        "curr_month": curr_month,
    }

    # Walk the slide sequence
    for slide_cfg in slides_config.get("slides", []):
        _dispatch_slide(prs, slide_cfg, layout_map, context)

    # Serialize
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


# ---------------------------------------------------------------------------
# Debug helper
# ---------------------------------------------------------------------------

def inspect_template(template_path: str) -> dict[str, Any]:
    """Enumerate all slides, shapes, and placeholders in a template."""
    path = Path(template_path)
    if not path.is_file():
        raise FileNotFoundError(f"Template file not found: {template_path}")

    prs = Presentation(str(path))
    result: dict[str, Any] = {
        "slide_count": len(prs.slides),
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
        "slides": [],
    }

    for slide_idx, slide in enumerate(prs.slides):
        slide_info: dict[str, Any] = {
            "slide_index": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": [],
        }

        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "has_text_frame": shape.has_text_frame,
            }

            if shape.is_placeholder:
                ph_fmt = shape.placeholder_format
                shape_info["is_placeholder"] = True
                shape_info["placeholder_idx"] = ph_fmt.idx
                shape_info["placeholder_type"] = str(ph_fmt.type)
            else:
                shape_info["is_placeholder"] = False

            if shape.has_text_frame:
                shape_info["text_preview"] = shape.text_frame.text[:200]

            slide_info["shapes"].append(shape_info)

        result["slides"].append(slide_info)

    return result
