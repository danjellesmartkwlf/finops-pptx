"""Shared PowerPoint helpers for the FinOps Report Generator.

Extracted from ``pptx_gen.py`` so that formatting and layout code is
separate from the config-driven slide orchestration logic.
"""

from __future__ import annotations

import calendar
import logging
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

TEMPLATE_CANDIDATES: list[str] = [
    "pptx_template/template1.pptx"
]

LAYOUT_TITLE_SLIDE = "Title Slide A"
LAYOUT_TRANSITION = "Transition Slide B"
LAYOUT_CONTENT = "Title, Body"
LAYOUT_TITLE_ONLY = "Title Only"

# Maximum data rows (excluding header/total) before splitting to a new slide
MAX_TABLE_ROWS_PER_SLIDE = 16

# "Dark Style 1 - Accent 4" built-in table style GUID
TABLE_STYLE_GUID = "{E929F9F4-4A8F-4326-A1B4-22849713DDAB}"


# ---------------------------------------------------------------------------
# Template resolution
# ---------------------------------------------------------------------------

def resolve_template_path(
    config: dict[str, Any],
    project_root: Path | None = None,
) -> Path:
    """Locate the .pptx template file on disk."""
    if project_root is None:
        project_root = Path(__file__).resolve().parent.parent

    for candidate in TEMPLATE_CANDIDATES:
        candidate_path = project_root / candidate
        if candidate_path.is_file():
            logger.info("Using template at %s", candidate_path)
            return candidate_path

    config_path_str: str = config.get("pptx", {}).get(
        "template_path", "assets/template.pptx"
    )
    config_path = project_root / config_path_str
    if config_path.is_file():
        logger.info("Using config-defined template at %s", config_path)
        return config_path

    searched = [
        str(project_root / c) for c in TEMPLATE_CANDIDATES
    ] + [str(config_path)]
    raise FileNotFoundError(
        f"PowerPoint template not found. Searched: {searched}"
    )


# ---------------------------------------------------------------------------
# Layout & slide helpers
# ---------------------------------------------------------------------------

def get_layout_by_name(prs: Any, name: str) -> Any:
    """Return a slide layout by its name, or raise ValueError."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    available = [layout.name for layout in prs.slide_layouts]
    raise ValueError(f"Layout '{name}' not found. Available: {available}")


def remove_all_slides(prs: Any) -> None:
    """Remove every slide from the presentation, keeping layouts intact."""
    sld_id_lst = prs.slides._sldIdLst
    while len(sld_id_lst):
        rId = sld_id_lst[0].rId
        prs.part.drop_rel(rId)
        del sld_id_lst[0]


# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

def parse_bullet_lines(text: str) -> list[dict[str, Any]]:
    """Parse narrative text that uses bullet markup.

    Convention (driven by config.yaml templates):
        ``- Text``       -> level 0, bold (section header)
        ``  - Text``     -> level 1, normal (detail item)
        empty line       -> blank spacer paragraph

    Plain text without any ``- `` prefixes falls through as level-0,
    non-bold paragraphs so that title/subtitle text is unaffected.
    """
    lines = text.split("\n")
    has_bullets = any(
        line.lstrip().startswith("- ") for line in lines if line.strip()
    )
    if not has_bullets:
        # Plain text path -- no bullet parsing
        return [{"text": line, "level": 0, "bold": False} for line in lines]

    items: list[dict[str, Any]] = []
    for line in lines:
        stripped = line.rstrip()
        if not stripped:
            # Blank line -> spacer paragraph
            items.append({"text": "", "level": 0, "bold": False})
            continue

        content = stripped.lstrip()
        indent = len(stripped) - len(content)

        if content.startswith("- "):
            content = content[2:]
            level = 1 if indent >= 2 else 0
            items.append({"text": content, "level": level, "bold": level == 0})
        else:
            items.append({"text": stripped, "level": 0, "bold": False})

    return items


def set_placeholder_text(
    placeholder: Any,
    text: str,
    font_size_pt: float | None = None,
) -> None:
    """Replace placeholder text while preserving the template's formatting.

    If *text* contains bullet markup (``- `` prefixes), paragraphs are
    assigned the appropriate ``level`` (0 or 1) and section headers are
    bolded.  Plain text is rendered as simple paragraphs.
    """
    tf = placeholder.text_frame

    # Capture existing formatting from the first run (if any)
    existing_font_name: str | None = None
    existing_font_size: Any = None
    existing_font_color: Any = None

    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        existing_font_name = first_run.font.name
        existing_font_size = first_run.font.size
        try:
            existing_font_color = first_run.font.color.rgb
        except AttributeError:
            existing_font_color = None

    tf.clear()

    items = parse_bullet_lines(text)
    for i, item in enumerate(items):
        if i == 0:
            paragraph = tf.paragraphs[0]
        else:
            paragraph = tf.add_paragraph()

        paragraph.level = item["level"]

        run = paragraph.add_run()
        run.text = item["text"]

        # Font formatting
        if existing_font_name:
            run.font.name = existing_font_name
        if font_size_pt is not None:
            run.font.size = Pt(font_size_pt)
        elif existing_font_size is not None:
            run.font.size = existing_font_size
        if item["bold"]:
            run.font.color.rgb = RGBColor(0xFF, 0x73, 0x0E)
        elif existing_font_color is not None:
            run.font.color.rgb = existing_font_color

        run.font.bold = item["bold"]


def set_text_by_idx(slide: Any, idx: int, text: str) -> bool:
    """Find a placeholder by idx on *slide* and set its text.

    Returns True if the placeholder was found and set, False otherwise.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            set_placeholder_text(shape, text)
            return True
    logger.warning("Placeholder idx %d not found on slide.", idx)
    return False


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def fmt_abbreviated(
    value: float | None,
    *,
    signed: bool = False,
    whole: bool = False,
) -> str:
    """Format a dollar value with abbreviated suffix (K, M, B).

    Args:
        value: The dollar amount. Returns "N/A" for None.
        signed: If True, prefix positive values with "+".
        whole: If True, round to whole numbers (e.g. "$5M" instead of "$5.81M").
    """
    if value is None:
        return "N/A"
    abs_val = abs(value)
    if signed:
        prefix = "+" if value > 0 else ("-" if value < 0 else "")
    else:
        prefix = "-" if value < 0 else ""
    dec = "0" if whole else "2"
    dec_k = "0" if whole else "1"
    if abs_val >= 1_000_000_000:
        return f"{prefix}${abs_val / 1_000_000_000:.{dec}f}B"
    if abs_val >= 1_000_000:
        return f"{prefix}${abs_val / 1_000_000:.{dec}f}M"
    if abs_val >= 1_000:
        return f"{prefix}${abs_val / 1_000:.{dec_k}f}K"
    return f"{prefix}${abs_val:,.0f}"


def fmt_abbreviated_signed(value: float | None) -> str:
    """Format with abbreviated suffix and explicit +/- sign."""
    return fmt_abbreviated(value, signed=True)


def fmt_abbreviated_whole(value: float | None) -> str:
    """Format with abbreviated suffix, rounded to whole numbers."""
    return fmt_abbreviated(value, whole=True)


def fmt_abbreviated_signed_whole(value: float | None) -> str:
    """Format with abbreviated suffix, explicit +/- sign, rounded to whole."""
    return fmt_abbreviated(value, signed=True, whole=True)


def fmt_daily_rate(value: float | None) -> str:
    """Format a dollar-per-day value (no sign prefix)."""
    if value is None:
        return "N/A"
    abs_val = abs(value)
    sign = "-" if value < 0 else ""
    if abs_val >= 1_000:
        return f"{sign}${abs_val / 1_000:.0f}K/day"
    return f"{sign}${abs_val:.0f}/day"


def fmt_daily_rate_signed(value: float | None) -> str:
    """Format a dollar-per-day change value with explicit +/- sign."""
    if value is None:
        return "N/A"
    abs_val = abs(value)
    sign = "+" if value > 0 else ("-" if value < 0 else "")
    if abs_val >= 1_000:
        return f"{sign}${abs_val / 1_000:.0f}K/day"
    return f"{sign}${abs_val:.0f}/day"


def fmt_obs_abbreviated(value: float | int | None) -> str:
    """Format an observation count with abbreviated suffix (B, T)."""
    if value is None:
        return "N/A"
    v = float(value)
    abs_val = abs(v)
    sign = "-" if v < 0 else ""
    if abs_val >= 1e12:
        return f"{sign}{abs_val / 1e12:.2f}T"
    if abs_val >= 1e9:
        return f"{sign}{abs_val / 1e9:.1f}B"
    if abs_val >= 1e6:
        return f"{sign}{abs_val / 1e6:.0f}M"
    return f"{sign}{abs_val:,.0f}"


def fmt_unit_cost(value: float | None) -> str:
    """Format a cost-per-1M value as '$X.XX'."""
    if value is None:
        return "N/A"
    return f"${value:.2f}"


def fmt_pct(value: float | None) -> str:
    """Format a percentage value with sign."""
    if value is None:
        return "N/A"
    sign = "+" if value > 0 else ""
    return f"{sign}{value:.1f}%"


def apply_table_style(table: Any) -> None:
    """Apply the 'Dark Style 1 - Accent 4' built-in style via GUID."""
    tbl_pr = table._tbl.tblPr

    # Remove any existing tableStyleId element (python-pptx default)
    for child in list(tbl_pr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("tableStyleId", "tblStyle"):
            tbl_pr.remove(child)

    style_el = etree.SubElement(tbl_pr, qn("a:tableStyleId"))
    style_el.text = TABLE_STYLE_GUID


def set_cell(
    cell: Any,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    alignment: Any = None,
    color: RGBColor | None = None,
) -> None:
    """Write text into a table cell with formatting."""
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    # Reduce internal margins for tighter cell spacing
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)

    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment

    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def font_size_for_row_count(data_row_count: int) -> int:
    """Return the appropriate font size based on data row count.

    - Under 10 data rows: 18pt
    - 10-14 data rows: 16pt
    - 15+ data rows: 14pt
    """
    if data_row_count >= 15:
        return 14
    if data_row_count >= 10:
        return 16
    return 18


def build_table_on_slide(
    slide: Any,
    headers: list[str],
    data_rows: list[list[str]],
    total_row: list[str] | None = None,
    col_widths: list[float] | None = None,
    font_size: int | None = None,
    top_inches: float = 1.4,
    row_height_factor: float = 0.32,
) -> tuple[Any, float]:
    """Create a styled table on a slide.

    Args:
        slide: The slide to add the table to.
        headers: Column header strings.
        data_rows: List of lists, each inner list is one row of cell strings.
        total_row: Optional totals row (same length as headers).
        col_widths: Optional column widths in inches.
        font_size: Optional font size override. If None, auto-computed
            from the data row count.
        top_inches: Top position in inches (default 1.4).
        row_height_factor: Height per row in inches (default 0.32).

    Returns:
        A (table, bottom_inches) tuple — the python-pptx Table object and
        the bottom edge position in inches.
    """
    n_rows = 1 + len(data_rows) + (1 if total_row else 0)
    n_cols = len(headers)

    if font_size is None:
        font_size = font_size_for_row_count(len(data_rows))

    left = Inches(0.5)
    top = Inches(top_inches)
    width = Inches(12.3)
    height = Inches(min(row_height_factor * n_rows, 5.8))

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    # Set tblPr flags (no banded columns)
    tbl_pr = table._tbl.tblPr
    tbl_pr.set("firstRow", "1")
    tbl_pr.set("lastRow", "1" if total_row else "0")
    tbl_pr.set("bandRow", "1")
    tbl_pr.set("bandCol", "0")
    tbl_pr.set("firstCol", "1")

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # -- Populate cells --
    for col_idx, header in enumerate(headers):
        align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        set_cell(table.cell(0, col_idx), header, font_size=font_size, bold=True, alignment=align)

    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            set_cell(table.cell(row_idx, col_idx), value, font_size=font_size, alignment=align)

    if total_row:
        t_row_idx = 1 + len(data_rows)
        for col_idx, value in enumerate(total_row):
            align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            set_cell(
                table.cell(t_row_idx, col_idx), value, font_size=font_size, bold=True, alignment=align
            )

    # -- Apply Dark Style 1 - Accent 4 via built-in GUID --
    apply_table_style(table)

    bottom_inches = top_inches + min(row_height_factor * n_rows, 5.8)
    return table, bottom_inches


def add_paginated_table_slides(
    prs: Any,
    layout: Any,
    title: str,
    headers: list[str],
    data_rows: list[list[str]],
    total_row: list[str] | None = None,
    col_widths: list[float] | None = None,
    font_size: int | None = None,
) -> None:
    """Add one or more slides with a paginated table.

    If data_rows exceeds ``MAX_TABLE_ROWS_PER_SLIDE``, splits across
    multiple slides.  The total row only appears on the last slide.

    Args:
        font_size: Optional font size override in points. If None,
            auto-computed from the per-page row count.
    """
    for page_start in range(0, max(len(data_rows), 1), MAX_TABLE_ROWS_PER_SLIDE):
        page_rows = data_rows[page_start : page_start + MAX_TABLE_ROWS_PER_SLIDE]
        is_last_page = (
            page_start + MAX_TABLE_ROWS_PER_SLIDE >= len(data_rows)
        )

        slide = prs.slides.add_slide(layout)
        set_text_by_idx(slide, 0, title)

        # Use explicit font_size if provided, otherwise auto-compute
        page_font_size = font_size if font_size is not None else font_size_for_row_count(len(page_rows))

        build_table_on_slide(
            slide,
            headers,
            page_rows,
            total_row=total_row if is_last_page else None,
            col_widths=col_widths,
            font_size=page_font_size,
        )


def compute_mom_headers(month_label: str) -> tuple[str, str]:
    """Compute column headers for the MoM table.

    Given month_label like "January 2026", returns abbreviated labels for
    the previous month and the current month: ("Dec 2025", "Jan 2026").
    """
    parts = month_label.split()
    month_name, year = parts[0], int(parts[1])
    month_num = list(calendar.month_name).index(month_name)

    # Current month label
    current_label = f"{calendar.month_abbr[month_num]} {year}"

    # Previous month label
    prev_m, prev_y = month_num - 1, year
    if prev_m == 0:
        prev_m, prev_y = 12, year - 1
    prev_label = f"{calendar.month_abbr[prev_m]} {prev_y}"

    return prev_label, current_label
